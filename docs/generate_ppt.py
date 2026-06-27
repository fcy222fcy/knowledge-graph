from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_flowchart_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 定义所有功能流程
    flows = [
        {
            "title": "5.1.1 注册功能实现流程",
            "steps": [
                {"text": "开始注册", "type": "start"},
                {"text": "填写注册信息\n用户名/邮箱/密码/确认密码", "type": "process"},
                {"text": "前端表单验证", "type": "decision"},
                {"text": "验证通过", "type": "arrow"},
                {"text": "发送POST请求\n/api/v1/auth/register", "type": "process"},
                {"text": "Controller接收请求\n绑定注册DTO", "type": "process"},
                {"text": "参数验证", "type": "decision"},
                {"text": "Service层处理", "type": "process"},
                {"text": "检查用户名/邮箱\n是否已存在", "type": "decision"},
                {"text": "bcrypt密码哈希", "type": "process"},
                {"text": "创建User实体\n写入数据库", "type": "process"},
                {"text": "返回成功\n跳转登录页", "type": "end"},
            ]
        },
        {
            "title": "5.1.2 登录功能实现流程",
            "steps": [
                {"text": "开始登录", "type": "start"},
                {"text": "输入用户名和密码", "type": "process"},
                {"text": "发送POST请求\n/api/v1/auth/login", "type": "process"},
                {"text": "查询用户记录", "type": "process"},
                {"text": "用户存在?", "type": "decision"},
                {"text": "bcrypt比对密码", "type": "process"},
                {"text": "密码正确?", "type": "decision"},
                {"text": "生成JWT Token\n编码用户信息", "type": "process"},
                {"text": "返回Token和用户信息", "type": "process"},
                {"text": "前端存储Token\nPinia + localStorage", "type": "process"},
                {"text": "Axios拦截器自动携带", "type": "process"},
                {"text": "路由守卫验证\n允许访问", "type": "end"},
            ]
        },
        {
            "title": "5.1.3 知识问答功能实现流程",
            "steps": [
                {"text": "开始问答", "type": "start"},
                {"text": "输入问题内容", "type": "process"},
                {"text": "发送POST请求\n/api/v1/ask", "type": "process"},
                {"text": "会话ID为空?", "type": "decision"},
                {"text": "创建新AskSession\n标题=问题前20字符", "type": "process"},
                {"text": "保存用户问题\nAskMessage role=user", "type": "process"},
                {"text": "知识检索\n关键词模糊匹配", "type": "process"},
                {"text": "计算置信度评分", "type": "process"},
                {"text": "生成回答内容", "type": "process"},
                {"text": "保存助手回答\nAskMessage role=assistant", "type": "process"},
                {"text": "返回回答+置信度\n+知识点列表", "type": "end"},
            ]
        },
        {
            "title": "5.1.4 知识图谱浏览功能实现流程",
            "steps": [
                {"text": "进入图谱页面", "type": "start"},
                {"text": "发送GET请求\n/api/v1/graph", "type": "process"},
                {"text": "查询KnowledgePoint表", "type": "process"},
                {"text": "查询KnowledgeRelation表", "type": "process"},
                {"text": "组装JSON响应\nnodes + edges", "type": "process"},
                {"text": "返回图谱数据", "type": "process"},
                {"text": "D3.js forceSimulation\n创建力导向图", "type": "process"},
                {"text": "设置力:\ncharge + link + center", "type": "process"},
                {"text": "SVG渲染\n节点=圆形\n关系=连线", "type": "process"},
                {"text": "用户交互:\n拖拽+缩放", "type": "end"},
            ]
        },
        {
            "title": "5.1.5 题库练习功能实现流程",
            "steps": [
                {"text": "进入题库页面", "type": "start"},
                {"text": "发送GET请求\n/api/v1/question", "type": "process"},
                {"text": "查询Question表\n返回题目列表", "type": "process"},
                {"text": "展示题目卡片\n标题/类型/难度", "type": "process"},
                {"text": "学生选择题目", "type": "process"},
                {"text": "发送POST请求\n/api/v1/quiz", "type": "process"},
                {"text": "查询正确答案", "type": "process"},
                {"text": "比对答案", "type": "decision"},
                {"text": "标记正确/错误", "type": "process"},
                {"text": "保存Quiz记录", "type": "process"},
                {"text": "返回结果+解析", "type": "end"},
            ]
        },
        {
            "title": "5.1.6 学习分析功能实现流程",
            "steps": [
                {"text": "进入学习分析页面", "type": "start"},
                {"text": "发送GET请求\n/api/v1/analytics/overview", "type": "process"},
                {"text": "查询AskSession表\n统计会话总数", "type": "process"},
                {"text": "查询AskMessage表\n统计消息总数", "type": "process"},
                {"text": "查询Quiz表\n统计答题正确数", "type": "process"},
                {"text": "按难度统计正确率\neasy/medium/hard", "type": "process"},
                {"text": "分析薄弱知识点", "type": "process"},
                {"text": "封装统计数据", "type": "process"},
                {"text": "返回统计数据", "type": "process"},
                {"text": "ECharts渲染图表\n问答+答题+雷达图", "type": "end"},
            ]
        },
        {
            "title": "5.2.1 教师登录功能实现流程",
            "steps": [
                {"text": "教师开始登录", "type": "start"},
                {"text": "输入教师账号和密码", "type": "process"},
                {"text": "发送POST请求\n/api/v1/teacher_auth/login", "type": "process"},
                {"text": "查询Teacher表", "type": "process"},
                {"text": "教师存在?", "type": "decision"},
                {"text": "bcrypt比对密码", "type": "process"},
                {"text": "密码正确?", "type": "decision"},
                {"text": "生成JWT Token\n包含教师身份标识", "type": "process"},
                {"text": "RequireTeacherAuth\n中间件验证", "type": "process"},
                {"text": "返回Token和教师信息", "type": "process"},
                {"text": "进入教师管理界面", "type": "end"},
            ]
        },
        {
            "title": "5.2.2 知识图谱构建功能实现流程",
            "steps": [
                {"text": "教师点击构建图谱", "type": "start"},
                {"text": "弹出资料选择对话框", "type": "process"},
                {"text": "勾选需要的文档", "type": "process"},
                {"text": "发送POST请求\n/api/v1/graph/build", "type": "process"},
                {"text": "查询Document记录\n获取解析内容", "type": "process"},
                {"text": "分块处理\n提取知识点和关系", "type": "process"},
                {"text": "确定性规则抽取", "type": "process"},
                {"text": "去重检查", "type": "decision"},
                {"text": "写入KnowledgePoint表", "type": "process"},
                {"text": "写入KnowledgeRelation表", "type": "process"},
                {"text": "记录构建统计", "type": "process"},
                {"text": "返回结果\n刷新图谱可视化", "type": "end"},
            ]
        },
    ]

    # 颜色定义
    BLACK = RGBColor(0, 0, 0)
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GRAY = RGBColor(240, 240, 240)
    DARK_GRAY = RGBColor(100, 100, 100)

    for flow in flows:
        slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)

        # 设置背景为白色
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = WHITE

        # 添加标题
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        title_frame = title_shape.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = flow["title"]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = BLACK
        p.alignment = PP_ALIGN.CENTER

        # 添加分隔线
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.1), Inches(12.333), Inches(0.02)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = BLACK
        line.line.fill.background()

        # 绘制流程图
        steps = flow["steps"]
        num_steps = len(steps)

        # 计算布局 - 分两列显示
        if num_steps <= 6:
            cols = 1
            rows = num_steps
        else:
            cols = 2
            rows = (num_steps + 1) // 2

        # 每个元素的尺寸
        box_width = Inches(2.8)
        box_height = Inches(0.6)
        h_gap = Inches(0.3)
        v_gap = Inches(0.15)

        # 起始位置
        start_x = Inches(0.8)
        start_y = Inches(1.4)

        for idx, step in enumerate(steps):
            # 计算当前位置
            if cols == 1:
                col = 0
                row = idx
            else:
                col = idx // rows
                row = idx % rows

            x = start_x + col * (box_width + Inches(3.5))
            y = start_y + row * (box_height + v_gap)

            # 根据类型选择形状和颜色
            if step["type"] == "start":
                shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
                fill_color = RGBColor(50, 50, 50)
                font_color = WHITE
            elif step["type"] == "end":
                shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
                fill_color = RGBColor(50, 50, 50)
                font_color = WHITE
            elif step["type"] == "decision":
                shape_type = MSO_SHAPE.DIAMOND
                fill_color = LIGHT_GRAY
                font_color = BLACK
                box_height_decision = Inches(0.8)
            elif step["type"] == "arrow":
                shape_type = MSO_SHAPE.RIGHT_ARROW
                fill_color = RGBColor(180, 180, 180)
                font_color = BLACK
            else:
                shape_type = MSO_SHAPE.RECTANGLE
                fill_color = WHITE
                font_color = BLACK

            # 创建形状
            if step["type"] == "decision":
                shape = slide.shapes.add_shape(shape_type, x, y, box_width, box_height_decision)
            else:
                shape = slide.shapes.add_shape(shape_type, x, y, box_width, box_height)

            # 设置样式
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
            shape.line.color.rgb = BLACK
            shape.line.width = Pt(1.5)

            # 添加文本
            text_frame = shape.text_frame
            text_frame.word_wrap = True
            text_frame.auto_size = None

            # 清除默认段落
            for paragraph in text_frame.paragraphs:
                paragraph.clear()

            # 添加新文本
            p = text_frame.paragraphs[0]
            p.text = step["text"]
            p.font.size = Pt(11)
            p.font.color.rgb = font_color
            p.alignment = PP_ALIGN.CENTER
            text_frame.paragraphs[0].font.bold = False

            # 垂直居中
            text_frame.paragraphs[0].space_before = Pt(4)

            # 添加连接箭头（除了最后一个）
            if idx < num_steps - 1:
                if cols == 1 or row < rows - 1:
                    # 向下箭头
                    arrow_x = x + box_width / 2 - Inches(0.1)
                    arrow_y = y + box_height + Inches(0.02)
                    arrow = slide.shapes.add_shape(
                        MSO_SHAPE.DOWN_ARROW,
                        arrow_x, arrow_y, Inches(0.2), Inches(0.12)
                    )
                    arrow.fill.solid()
                    arrow.fill.fore_color.rgb = BLACK
                    arrow.line.fill.background()

        # 添加页码
        page_num = slide.shapes.add_textbox(Inches(12.5), Inches(7.0), Inches(0.8), Inches(0.4))
        page_frame = page_num.text_frame
        p = page_frame.paragraphs[0]
        p.text = f"{flows.index(flow) + 1}/8"
        p.font.size = Pt(10)
        p.font.color.rgb = DARK_GRAY
        p.alignment = PP_ALIGN.RIGHT

    # 保存文件
    output_path = r"e:\goCode\goFile\software engineering\docs\系统功能实现流程图.pptx"
    prs.save(output_path)
    print(f"PPT已生成: {output_path}")

if __name__ == "__main__":
    create_flowchart_ppt()
