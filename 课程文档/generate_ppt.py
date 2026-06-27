# -*- coding: utf-8 -*-
"""
SE智图问答平台 答辩PPT生成脚本
基于现有模板风格：16:9, 蓝灰色系, 卡片式布局
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# ── 颜色常量（对齐模板） ──
DARK      = RGBColor(0x1F, 0x29, 0x37)  # 主标题
SUBTITLE  = RGBColor(0x1F, 0x23, 0x29)  # 副标题
BLUE      = RGBColor(0x3B, 0x82, 0xF6)  # 蓝色强调
GRAY      = RGBColor(0x6B, 0x72, 0x80)  # 正文灰
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF3, 0xF4, 0xF6)  # 浅灰背景卡片
CARD_BG   = RGBColor(0xFF, 0xFF, 0xFF)  # 白色卡片

# ── EMU 常量 ──
SLIDE_W = 12192000
SLIDE_H = 6858000
MARGIN_L = 1016000   # 左边距 ~1英寸
MARGIN_T = 762000    # 上边距 ~0.8英寸


def set_slide_bg(slide, color):
    """设置幻灯片背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    """添加矩形形状"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()  # 默认无边框
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width or Emu(12700)
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color=None):
    """添加圆角矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=Pt(11), bold=False,
             color=DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font_name='微软雅黑'):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, font_size=Pt(11),
                       bold=False, color=DARK, line_spacing=Pt(18), align=PP_ALIGN.LEFT):
    """添加多行文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = '微软雅黑'
        p.alignment = align
        p.space_before = Pt(2)
        p.space_after = Pt(2)
        if line_spacing:
            p.line_spacing = line_spacing
    return txBox


def add_card(slide, left, top, width, height, title, desc, accent_color=BLUE,
             title_size=Pt(16), desc_size=Pt(11), card_color=CARD_BG):
    """添加带顶部色条的卡片"""
    # 卡片背景
    card = add_rounded_rect(slide, left, top, width, height, card_color)
    # 顶部色条
    add_shape(slide, left, top, width, Emu(60960), accent_color)
    # 标题
    add_text(slide, left + Emu(152400), top + Emu(203200),
             width - Emu(304800), Emu(406400),
             title, title_size, bold=True, color=DARK)
    # 描述
    add_text(slide, left + Emu(152400), top + Emu(635000),
             width - Emu(304800), height - Emu(762000),
             desc, desc_size, color=GRAY)
    return card


def add_page_title(slide, title):
    """添加页面标题"""
    add_text(slide, MARGIN_L, MARGIN_T, Emu(10160000), Emu(508000),
             title, Pt(32), bold=True, color=DARK)


def add_bullet_card(slide, left, top, width, height, title, bullets, accent_color=BLUE):
    """添加带项目列表的卡片"""
    card = add_rounded_rect(slide, left, top, width, height, CARD_BG)
    add_shape(slide, left, top, width, Emu(60960), accent_color)
    add_text(slide, left + Emu(152400), top + Emu(203200),
             width - Emu(304800), Emu(406400),
             title, Pt(16), bold=True, color=DARK)
    lines = ['• ' + b for b in bullets]
    add_multiline_text(slide, left + Emu(152400), top + Emu(635000),
                       width - Emu(304800), height - Emu(762000),
                       lines, Pt(11), color=GRAY, line_spacing=Pt(20))


# ============================================================
# 创建演示文稿
# ============================================================
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ============================================================
# 第1页：封面
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
set_slide_bg(slide1, WHITE)

# 背景装饰大矩形
add_shape(slide1, 0, 0, SLIDE_W, SLIDE_H, LIGHT_BG)

# 左侧蓝色竖条装饰
add_shape(slide1, 0, 0, Emu(304800), SLIDE_H, BLUE)

# 项目名称
add_text(slide1, Emu(1270000), Emu(1905000), Emu(9652000), Emu(889000),
         'SE智图问答平台', Pt(48), bold=True, color=DARK, align=PP_ALIGN.LEFT)

# 副标题
add_text(slide1, Emu(1270000), Emu(2921000), Emu(9652000), Emu(508000),
         '基于知识图谱的软件工程课程智能问答系统', Pt(24), color=SUBTITLE)

# 蓝色分割线
add_shape(slide1, Emu(1270000), Emu(3556000), Emu(2032000), Emu(38100), BLUE)

# 英文标识
add_text(slide1, Emu(1270000), Emu(3810000), Emu(9652000), Emu(381000),
         'INTELLIGENT Q&A SYSTEM · 2026', Pt(12), color=GRAY)

# 成员分工表格区域
add_text(slide1, Emu(1270000), Emu(4445000), Emu(3048000), Emu(381000),
         '小组成员与分工', Pt(14), bold=True, color=DARK)

# 成员信息
members = [
    ('符成岩（组长）', '后端架构、数据库设计、RAG 问答、部署脚本'),
    ('芮晓雨', 'Vue 3 学生端 + 教师端、D3.js 图谱可视化、ECharts 图表'),
    ('苗琳', 'Python 知识抽取服务、黑盒测试、课程文档'),
]
for i, (name, role) in enumerate(members):
    y = Emu(4953000 + i * 457200)
    add_text(slide1, Emu(1270000), y, Emu(2286000), Emu(381000),
             name, Pt(12), bold=True, color=DARK)
    add_text(slide1, Emu(3683000), y, Emu(7239000), Emu(381000),
             role, Pt(11), color=GRAY)

# 主要技术栈
add_text(slide1, Emu(1270000), Emu(6350000), Emu(9652000), Emu(304800),
         '技术栈：Vue 3 + Go Gin + MySQL + Neo4j + Ollama（Qwen3:8b）', Pt(11), color=GRAY)


# ============================================================
# 第2页：背景与需求
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2, WHITE)

# 页面标题
add_text(slide2, MARGIN_L, MARGIN_T, Emu(10160000), Emu(508000),
         '项目背景与目标', Pt(32), bold=True, color=DARK)

# 左侧区域 - 背景痛点
add_text(slide2, MARGIN_L, Emu(1524000), Emu(4572000), Emu(381000),
         '背景痛点', Pt(18), bold=True, color=BLUE)

# 4个痛点卡片（2x2网格）
pains = [
    ('知识获取低效', '海量课程资料分散无序，学生难以快速定位知识点'),
    ('知识体系零散', '知识点间关联被隐藏，无法构建系统化网络'),
    ('评估方式单一', '传统考试难以动态、个性化评估学习效果'),
    ('资源更新滞后', '教师手动更新耗时，内容滞后于技术发展'),
]
for i, (t, d) in enumerate(pains):
    col = i % 2
    row = i // 2
    x = Emu(1016000 + col * 2540000)
    y = Emu(2032000 + row * 2159000)
    add_card(slide2, x, y, Emu(2413000), Emu(1905000), t, d)

# 右侧区域 - 项目目标
add_text(slide2, Emu(6223000), Emu(1524000), Emu(4953000), Emu(381000),
         '项目目标', Pt(18), bold=True, color=BLUE)

goals = [
    ('知识图谱可视化', '将课程文档转化为知识图谱，结构化存储与可视化展示'),
    ('智能问答服务', '基于 RAG + 本地 LLM，提供精准自然语言问答'),
    ('学习分析评估', '分析答题记录，自动评估知识点掌握度，生成可视化报告'),
]
for i, (t, d) in enumerate(goals):
    y = Emu(2032000 + i * 1651000)
    add_card(slide2, Emu(6223000), y, Emu(4953000), Emu(1524000), t, d)

# 底部角色说明
add_text(slide2, MARGIN_L, Emu(6477000), Emu(10160000), Emu(304800),
         '两类角色：学生端（问答、图谱浏览、练习、分析）+ 教师端（图谱构建、题库管理、资料审核）',
         Pt(11), color=GRAY, align=PP_ALIGN.LEFT)


# ============================================================
# 第3页：项目成员组成与分工
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3, WHITE)
add_page_title(slide3, '项目成员组成与分工')

members = [
    ('符成岩（组长）', '20231514301', '后端开发与系统架构',
     ['系统架构设计（Go Gin + Vue 3 + MySQL + Neo4j + MinIO）',
      '后端核心开发：用户认证、资料管理、知识图谱、智能问答、题库、学习分析',
      '数据库设计：MySQL 业务表 + Neo4j 图谱模型',
      '应用生命周期管理与部署脚本（start.bat / stop.bat）']),
    ('芮晓雨', '——', '前端开发',
     ['学生端 Vue 3 SPA：登录注册、问答、图谱、练习、统计等全部页面',
      '教师端 admin-vue-app：管理仪表盘、资料审核、题目/学生管理',
      'D3.js 知识图谱力导向布局可视化',
      'ECharts 学习统计图表、Vue Router 路由鉴权、Pinia 状态管理']),
    ('苗琳', '——', '测试、文档与 AI 模块',
     ['Python 知识抽取服务：实体识别、关系抽取、FAISS 向量化',
      '黑盒测试用例设计与执行（30+ 用例）',
      '全套课程文档编写（需求分析、设计文档、测试文档等）']),
]

for i, (name, sid, role, tasks) in enumerate(members):
    y_base = Emu(1651000 + i * 1778000)
    add_rounded_rect(slide3, Emu(1016000), y_base, Emu(10160000), Emu(1651000), CARD_BG)
    add_shape(slide3, Emu(1016000), y_base, Emu(76200), Emu(1651000), BLUE)
    add_text(slide3, Emu(1397000), y_base + Emu(152400), Emu(2540000), Emu(381000),
             name, Pt(20), bold=True, color=DARK)
    # 角色标签
    add_rounded_rect(slide3, Emu(4064000), y_base + Emu(177800), Emu(1778000), Emu(304800), BLUE)
    add_text(slide3, Emu(4064000), y_base + Emu(177800), Emu(1778000), Emu(304800),
             role, Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    task_lines = ['• ' + t for t in tasks]
    add_multiline_text(slide3, Emu(1397000), y_base + Emu(609600), Emu(9525000), Emu(1016000),
                       task_lines, Pt(11), color=GRAY, line_spacing=Pt(22))


# ============================================================
# 第4页：技术架构
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide4, WHITE)
add_page_title(slide4, '系统技术架构')

layers = [
    ('前端展示层 (Vue 3)', '用户交互与界面渲染核心',
     'Vue 3 + TypeScript + Element Plus + D3.js + ECharts，通过 HTTP/REST API 与后端通信'),
    ('后端服务层 (Go + Gin)', '业务逻辑与 API 网关中枢',
     'Gin RESTful API，MVC 分层：Controller → Service → Repository，JWT 鉴权'),
    ('AI 智能服务层 (Python)', '计算密集型任务处理单元',
     'Ollama 部署 Qwen3:8b + Nomic-embed-text，RAG 检索增强生成、LLM 知识抽取'),
    ('多维数据存储层', '混合架构存储解决方案',
     'MySQL（业务数据）+ Neo4j（知识图谱）+ MinIO（文件存储）'),
]
for i, (title, subtitle, desc) in enumerate(layers):
    col, row = i % 2, i // 2
    x = Emu(1016000 + col * 5461000)
    y = Emu(1651000 + row * 2540000)
    w, h = Emu(5080000), Emu(2286000)

    add_rounded_rect(slide4, x, y, w, h, CARD_BG)
    add_shape(slide4, x, y, Emu(76200), h, BLUE)
    icon_size = Emu(609600)
    add_rounded_rect(slide4, x + Emu(304800), y + Emu(254000), icon_size, icon_size, BLUE)
    add_text(slide4, x + Emu(1143000), y + Emu(254000), w - Emu(1447800), Emu(381000),
             title, Pt(18), bold=True, color=DARK)
    add_text(slide4, x + Emu(1143000), y + Emu(685800), w - Emu(1447800), Emu(254000),
             subtitle, Pt(12), color=BLUE)
    add_text(slide4, x + Emu(304800), y + Emu(1143000), w - Emu(609600), Emu(1016000),
             desc, Pt(11), color=GRAY)


# ============================================================
# 第5页：核心功能演示 - 智能问答 RAG
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide5, WHITE)
add_page_title(slide5, '核心功能 — 智能问答（RAG）')

# RAG 流程（横向5步骤）
steps = [
    ('用户提问', '自然语言输入\n支持多轮对话'),
    ('文本向量化', 'Nomic-embed-text\n768维向量'),
    ('向量检索', '余弦相似度\nTop-K匹配'),
    ('LLM生成', 'Qwen3:8b\n检索增强生成'),
    ('流式输出', 'SSE推送\n逐Token返回'),
]
for i, (title, desc) in enumerate(steps):
    x = Emu(609600 + i * 2286000)
    y = Emu(1651000)
    w, h = Emu(2032000), Emu(2032000)
    add_rounded_rect(slide5, x, y, w, h, CARD_BG)
    add_shape(slide5, x, y, w, Emu(60960), BLUE)
    num_size = Emu(508000)
    add_rounded_rect(slide5, x + Emu(762000), y + Emu(203200), num_size, num_size, BLUE)
    add_text(slide5, x + Emu(762000), y + Emu(228600), num_size, num_size,
             str(i + 1), Pt(24), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide5, x + Emu(152400), y + Emu(812800), w - Emu(304800), Emu(381000),
             title, Pt(16), bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text(slide5, x + Emu(152400), y + Emu(1270000), w - Emu(304800), Emu(609600),
             desc, Pt(11), color=GRAY, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        arrow_x = x + w + Emu(38100)
        add_text(slide5, arrow_x, y + Emu(812800), Emu(152400), Emu(381000),
                 '→', Pt(24), bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# 下方两栏
add_bullet_card(slide5, Emu(1016000), Emu(4064000), Emu(5080000), Emu(2413000),
    '降级方案', [
        'Ollama 不可用时自动切换为关键词匹配',
        '基于 KnowledgePoint 表模糊查询',
        '保证 AI 异常时仍可提供基本问答',
    ])
add_bullet_card(slide5, Emu(6223000), Emu(4064000), Emu(5080000), Emu(2413000),
    '多轮对话支持', [
        'AskSession 维护会话上下文',
        '历史消息作为 LLM 输入上下文',
        '置信度评分展示回答可靠程度',
    ])


# ============================================================
# 第6页：核心功能演示 - 知识图谱 & 题库
# ============================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide6, WHITE)
add_page_title(slide6, '核心功能 — 知识图谱 & 题库')

add_bullet_card(slide6, Emu(1016000), Emu(1651000), Emu(5080000), Emu(4826000),
    '知识图谱构建与浏览', [
        '教师选择已审核文档，一键触发自动构建',
        'LLM Prompt 引导抽取知识点与关系（JSON）',
        '降级方案：正则 + 软件工程本体关键词规则抽取',
        '去重检查 → 写入 Neo4j 图数据库',
        'D3.js 力导向布局可视化渲染',
        '支持节点拖拽、滚轮缩放、点击查看详情',
        '不同分类知识点使用不同颜色区分',
        '构建统计：新增知识点数、关系数、跳过重复数',
    ])

add_bullet_card(slide6, Emu(6223000), Emu(1651000), Emu(5080000), Emu(2413000),
    '题库练习', [
        '题目列表支持按难度、知识点筛选',
        '单选/多选题型，难度分 easy/medium/hard',
        '选择答案提交 → 即时评分 + 详细解析',
        '答题记录持久化，支持历史查看',
    ])

add_bullet_card(slide6, Emu(6223000), Emu(4241800), Emu(5080000), Emu(2235200),
    '学习分析 & 教师端', [
        'ECharts 统计图表：问答次数、正确率、薄弱知识点',
        '教师端独立管理后台（admin-vue-app）',
        '资料审核、知识点/题目 CRUD、学生管理',
    ])


# ============================================================
# 第7页：测试
# ============================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide7, WHITE)
add_page_title(slide7, '系统测试')

tests = [
    ('单元测试 & 端到端测试', [
        '后端 Service 层核心逻辑单元测试',
        'Repository 层数据库操作测试',
        '前端组件渲染与交互测试',
        'API 接口端到端联调测试',
    ]),
    ('黑盒测试', [
        '30+ 测试用例，覆盖全部核心功能',
        'Postman / curl 构造 HTTP 请求验证',
        '覆盖：注册登录、智能问答、图谱构建',
        '覆盖：题库练习、资料管理、教师端',
        '涵盖正常流程与异常边界条件',
    ]),
    ('自动化测试', [
        'start.bat / stop.bat 一键启停脚本',
        'GORM AutoMigrate 自动建表',
        '种子数据自动初始化（首次启动）',
        'Vite 热重载前端开发自动化',
    ]),
]

for i, (title, bullets) in enumerate(tests):
    x = Emu(1016000 + i * 3556000)
    w = Emu(3429000)
    h = Emu(4572000)
    add_bullet_card(slide7, x, Emu(1651000), w, h, title, bullets)

add_shape(slide7, Emu(1016000), Emu(6350000), Emu(10160000), Emu(304800), LIGHT_BG)
add_text(slide7, Emu(1016000), Emu(6350000), Emu(10160000), Emu(304800),
         '✅ 测试结论：各功能模块运行正常，系统整体稳定可靠',
         Pt(13), bold=True, color=DARK, align=PP_ALIGN.CENTER)


# ============================================================
# 第8页：总结 + 致谢
# ============================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide8, WHITE)

add_text(slide8, MARGIN_L, MARGIN_T, Emu(10160000), Emu(508000),
         '总结 — 成果与不足', Pt(32), bold=True, color=DARK)

# 左侧成果
add_text(slide8, Emu(1016000), Emu(1524000), Emu(4572000), Emu(381000),
         '项目成果', Pt(18), bold=True, color=BLUE)

achievements = [
    '完整前后端分离平台（Vue 3 + Go Gin）',
    '集成 RAG + 本地 LLM（Qwen3:8b）智能问答',
    'LLM 驱动的知识图谱自动构建，具备降级方案',
    'D3.js 知识图谱可视化 + ECharts 学习分析',
    '学生/教师双角色完整功能闭环',
    '30+ 黑盒测试用例全部通过',
]
for i, text in enumerate(achievements):
    y = Emu(2032000 + i * 609600)
    dot = slide8.shapes.add_shape(MSO_SHAPE.OVAL, Emu(1016000), y + Emu(60960),
                                   Emu(127000), Emu(127000))
    dot.fill.solid()
    dot.fill.fore_color.rgb = BLUE
    dot.line.fill.background()
    add_text(slide8, Emu(1397000), y, Emu(4267200), Emu(304800),
             text, Pt(12), color=DARK)

# 右侧不足
add_text(slide8, Emu(6223000), Emu(1524000), Emu(4953000), Emu(381000),
         '不足与展望', Pt(18), bold=True, color=BLUE)

improvements = [
    ('知识抽取精度有限', '可引入 Few-shot Learning、更强 NER/RE 模型'),
    ('缺少容器化部署', '可使用 Docker Compose 一键部署所有服务'),
    ('测试以黑盒为主', '可补充单元测试和集成测试，完善质量保障'),
    ('用户体验待优化', '响应式布局、错误提示友好度、交互细节'),
]
for i, (t, d) in enumerate(improvements):
    y = Emu(2032000 + i * 914400)
    add_card(slide8, Emu(6223000), y, Emu(4953000), Emu(838200), t, d,
             title_size=Pt(14), desc_size=Pt(10))

# 底部致谢
add_shape(slide8, 0, Emu(5715000), SLIDE_W, Emu(1143000), LIGHT_BG)
add_text(slide8, 0, Emu(5842400), SLIDE_W, Emu(609600),
         '感谢指导老师与各位评委', Pt(20), bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_text(slide8, 0, Emu(6350000), SLIDE_W, Emu(304800),
         'THANKS', Pt(14), color=GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# 保存文件
# ============================================================
output_path = '课程文档/SE智图问答平台-答辩PPT.pptx'
prs.save(output_path)
print(f'✅ PPT 已生成: {output_path}')
print(f'   共 {len(prs.slides)} 页')
